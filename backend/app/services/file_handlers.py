import os
import logging
import datetime
from io import BytesIO
from typing import Optional

from fastapi import HTTPException
from fastapi.responses import StreamingResponse
import pandas as pd
from PyPDF2 import PdfReader
from PyPDF2.errors import PdfReadError
import pytesseract
from PIL import Image, ImageDraw
from pdf2image import convert_from_bytes
from docx import Document
from pptx import Presentation

from app.services.redaction import redact_by_level
from app.database import execute_query

# Resolve backend base directory
BASE_DIR = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

# Configure Tesseract OCR binary path if present locally
local_tess = os.path.join(BASE_DIR, "tools", "tesseract", "tesseract.exe")
tess_env = os.environ.get("TESSERACT_CMD")
if os.path.exists(local_tess):
    pytesseract.pytesseract.tesseract_cmd = local_tess
elif tess_env and os.path.exists(tess_env):
    pytesseract.pytesseract.tesseract_cmd = tess_env
elif os.path.exists(r"C:\Program Files\Tesseract-OCR\tesseract.exe"):
    pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

def get_poppler_path() -> Optional[str]:
    """
    Resolve Poppler library path for PDF to image conversion.
    """
    local_poppler = os.path.join(BASE_DIR, "tools", "poppler", "bin")
    if os.path.exists(local_poppler):
        return local_poppler
    local_poppler_lib = os.path.join(BASE_DIR, "tools", "poppler", "Library", "bin")
    if os.path.exists(local_poppler_lib):
        return local_poppler_lib
    env_path = os.environ.get("POPPLER_PATH")
    if env_path and os.path.exists(env_path):
        return env_path
    common_paths = [
        r"C:\Program Files\poppler\Library\bin",
        r"C:\Program Files (x86)\poppler\Library\bin",
        r"C:\Users\User\Downloads\Release-24.08.0-0\poppler-24.08.0\Library\bin"
    ]
    for p in common_paths:
        if os.path.exists(p):
            return p
    return None

async def record_history(user_id: Optional[int], filename: str, operation_type: str, redaction_level: int, status: str, details: str = ""):
    """
    Log redaction operations asynchronously into SQLite history table.
    """
    try:
        timestamp = datetime.datetime.now().isoformat()
        await execute_query('''
            INSERT INTO history (user_id, filename, operation_type, redaction_level, status, timestamp, details)
            VALUES (?, ?, ?, ?, ?, ?, ?)
        ''', (user_id or 0, filename, operation_type, redaction_level, status, timestamp, details))
    except Exception as e:
        logging.warning(f"Failed to log operation history: {e}")

def redact_faces_in_image(image: Image.Image, method: str = "blur", redaction_level: int = 1) -> Image.Image:
    """
    Detects and de-identifies visual biometric PII based on redaction level:
    - Levels 1-2: No visual masking (OCR text redaction only)
    - Level 3: Human Faces only
    - Level 4: Human Faces + Handwritten Signatures
    - Level 5: Human Faces + Signatures + Official Colored Stamps & Fingerprints/Thumbprints
    """
    if redaction_level < 3:
        return image

    try:
        import cv2
        import numpy as np
        
        img_np = np.array(image)
        if len(img_np.shape) == 3 and img_np.shape[2] == 3:
            img_bgr = cv2.cvtColor(img_np, cv2.COLOR_RGB2BGR)
        else:
            img_bgr = img_np.copy()
            if len(img_bgr.shape) == 2:
                img_bgr = cv2.cvtColor(img_bgr, cv2.COLOR_GRAY2BGR)
                
        gray = cv2.cvtColor(img_bgr, cv2.COLOR_BGR2GRAY)
        h, w = gray.shape
        
        boxes_to_redact = []

        # 1. FACE DETECTION (Haar Cascade) - Runs at levels 3 to 5
        if redaction_level >= 3:
            try:
                cascade_path = cv2.data.haarcascades + 'haarcascade_frontalface_default.xml'
                face_cascade = cv2.CascadeClassifier(cascade_path)
                faces = face_cascade.detectMultiScale(gray, scaleFactor=1.1, minNeighbors=5, minSize=(30, 30))
                for (fx, fy, fw, fh) in faces:
                    boxes_to_redact.append((fx, fy, fw, fh, "face"))
            except Exception as e_face:
                logging.debug(f"Face detection skipped: {e_face}")

        # 2. COLORED STAMPS, SEALS & THUMBPRINTS - Runs only at Level 5
        if redaction_level >= 5:
            try:
                hsv = cv2.cvtColor(img_bgr, cv2.COLOR_BGR2HSV)
                sat_mask = hsv[:, :, 1] > 60
                val_mask = (hsv[:, :, 2] > 30) & (hsv[:, :, 2] < 240)
                color_mask = ((sat_mask & val_mask) * 255).astype(np.uint8)
                
                kernel = cv2.getStructuringElement(cv2.MORPH_ELLIPSE, (7, 7))
                closed_color = cv2.morphologyEx(color_mask, cv2.MORPH_CLOSE, kernel)
                contours, _ = cv2.findContours(closed_color, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
                
                for cnt in contours:
                    area = cv2.contourArea(cnt)
                    if 400 < area < (h * w * 0.4):
                        cx, cy, cw, ch = cv2.boundingRect(cnt)
                        aspect = cw / float(max(1, ch))
                        if 0.3 < aspect < 5.0:
                            boxes_to_redact.append((cx, cy, cw, ch, "stamp_thumbprint"))
            except Exception as e_color:
                logging.debug(f"Color stamp/thumbprint detection skipped: {e_color}")

        # 3. HANDWRITTEN SIGNATURES (Level 4+) & GRAYSCALE THUMBPRINTS (Level 5)
        if redaction_level >= 4:
            try:
                thresh = cv2.adaptiveThreshold(gray, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY_INV, 15, 8)
                sig_kernel = cv2.getStructuringElement(cv2.MORPH_RECT, (15, 5))
                connected_strokes = cv2.morphologyEx(thresh, cv2.MORPH_CLOSE, sig_kernel)
                
                contours, _ = cv2.findContours(connected_strokes, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
                for cnt in contours:
                    area = cv2.contourArea(cnt)
                    if 800 < area < (h * w * 0.25):
                        sx, sy, sw, sh = cv2.boundingRect(cnt)
                        aspect = sw / float(max(1, sh))
                        hull = cv2.convexHull(cnt)
                        hull_area = cv2.contourArea(hull)
                        solidity = float(area) / max(1.0, hull_area)
                        
                        # Cursive signatures: wide aspect ratio (>1.8) and low solidity (<0.6) -> Level 4+
                        if aspect > 1.8 and solidity < 0.6 and sw > 80:
                            boxes_to_redact.append((sx, sy, sw, sh, "signature"))
                        # Grayscale fingerprints: aspect ratio close to 1.0 (0.7-1.4) and high ridge solidity (>0.6) -> Level 5 only
                        elif redaction_level >= 5 and (0.7 < aspect < 1.4 and solidity > 0.6 and sh > 60 and sw > 60):
                            boxes_to_redact.append((sx, sy, sw, sh, "fingerprint"))
            except Exception as e_sig:
                logging.debug(f"Signature detection skipped: {e_sig}")

        # Apply irreversible de-identification to all visual PII regions
        for (bx, by, bw, bh, pii_type) in boxes_to_redact:
            pad_w = int(bw * 0.08)
            pad_h = int(bh * 0.08)
            x1 = max(0, bx - pad_w)
            y1 = max(0, by - pad_h)
            x2 = min(w, bx + bw + pad_w)
            y2 = min(h, by + bh + pad_h)
            
            if method == "blackout":
                cv2.rectangle(img_bgr, (x1, y1), (x2, y2), (0, 0, 0), -1)
            else: # blur and pixelate
                roi = img_bgr[y1:y2, x1:x2]
                if roi.size > 0:
                    blurred = cv2.GaussianBlur(roi, (89, 89), 30)
                    img_bgr[y1:y2, x1:x2] = blurred

        if len(img_bgr.shape) == 3 and img_bgr.shape[2] == 3:
            img_rgb = cv2.cvtColor(img_bgr, cv2.COLOR_BGR2RGB)
        else:
            img_rgb = cv2.cvtColor(img_bgr, cv2.COLOR_GRAY2RGB)
        return Image.fromarray(img_rgb)
    except Exception as e:
        logging.warning(f"Visual biometric PII redaction failed: {e}")
        return image

async def handle_pdf(file_data: bytes, filename: str, redaction_level: int, user_id: Optional[int] = None) -> StreamingResponse:
    try:
        pdf_file = BytesIO(file_data)
        try:
            reader = PdfReader(pdf_file)
            _ = reader.pages[0]
        except Exception:
            raise HTTPException(status_code=400, detail="Invalid or corrupted PDF file.")
        pdf_file.seek(0)

        poppler_path = get_poppler_path()
        try:
            images = convert_from_bytes(pdf_file.read(), poppler_path=poppler_path)
        except Exception as poppler_error:
            raise HTTPException(
                status_code=500,
                detail=f"Error converting PDF to images. Ensure Poppler is installed. {str(poppler_error)}"
            )

        redacted_images = []
        for image in images:
            if redaction_level == 0:
                redacted_images.append(image)
                continue

            # Automatically de-identify human faces and biometrics based on redaction level
            image = redact_faces_in_image(image, method="blur", redaction_level=redaction_level)

            # SINGLE OCR CALL: get word bounding boxes and tokens in one pass
            draw = ImageDraw.Draw(image)
            tsv_data = pytesseract.image_to_data(image, output_type=pytesseract.Output.DATAFRAME)
            
            # Reconstruct full line text to allow multi-word entity detection
            valid_rows = tsv_data[tsv_data['text'].notnull() & (tsv_data['text'].str.strip() != "")]
            full_page_text = " ".join(valid_rows['text'].astype(str).tolist())
            redacted_page = redact_by_level(full_page_text, redaction_level) if full_page_text else ""

            for _, row in valid_rows.iterrows():
                text = str(row["text"]).strip()
                if text:
                    redacted_word = redact_by_level(text, redaction_level)
                    # Redact if the word itself is redacted OR if it was redacted within the full page context
                    if redacted_word != text or (len(text) > 2 and text not in redacted_page):
                        left, top, width, height = int(row['left']), int(row['top']), int(row['width']), int(row['height'])
                        draw.rectangle([left, top, left + width, top + height], fill="black")

            redacted_images.append(image)

        output_pdf = BytesIO()
        redacted_images[0].save(
            output_pdf,
            format="PDF",
            save_all=True,
            append_images=redacted_images[1:]
        )
        output_pdf.seek(0)

        await record_history(user_id, filename, "PDF Redaction", redaction_level, "SUCCESS")

        return StreamingResponse(
            output_pdf,
            media_type="application/pdf",
            headers={"Content-Disposition": f"attachment; filename=redacted_{filename}"}
        )
    except HTTPException as http_err:
        await record_history(user_id, filename, "PDF Redaction", redaction_level, "FAILED", str(http_err.detail))
        raise http_err
    except Exception as e:
        await record_history(user_id, filename, "PDF Redaction", redaction_level, "FAILED", str(e))
        raise HTTPException(status_code=500, detail=f"Error processing PDF: {str(e)}")

async def handle_txt(file_data: bytes, filename: str, redaction_level: int, user_id: Optional[int] = None) -> StreamingResponse:
    try:
        text = file_data.decode("utf-8", errors="replace")
        text = redact_by_level(text, redaction_level)
        output_file = BytesIO(text.encode("utf-8"))

        await record_history(user_id, filename, "TXT Redaction", redaction_level, "SUCCESS")

        return StreamingResponse(
            output_file,
            media_type="text/plain",
            headers={"Content-Disposition": f"attachment; filename=redacted_{filename}"}
        )
    except Exception as e:
        await record_history(user_id, filename, "TXT Redaction", redaction_level, "FAILED", str(e))
        raise HTTPException(status_code=500, detail=f"Error processing TXT file: {str(e)}")

async def handle_docx(file_data: bytes, filename: str, redaction_level: int, user_id: Optional[int] = None) -> StreamingResponse:
    try:
        doc = Document(BytesIO(file_data))
        for paragraph in doc.paragraphs:
            paragraph.text = redact_by_level(paragraph.text, redaction_level)

        output_docx = BytesIO()
        doc.save(output_docx)
        output_docx.seek(0)

        await record_history(user_id, filename, "DOCX Redaction", redaction_level, "SUCCESS")

        return StreamingResponse(
            output_docx,
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            headers={"Content-Disposition": f"attachment; filename=redacted_{filename}"}
        )
    except Exception as e:
        await record_history(user_id, filename, "DOCX Redaction", redaction_level, "FAILED", str(e))
        raise HTTPException(status_code=500, detail=f"Error processing DOCX file: {str(e)}")

async def handle_pptx(file_data: bytes, filename: str, redaction_level: int, user_id: Optional[int] = None) -> StreamingResponse:
    try:
        presentation = Presentation(BytesIO(file_data))
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.text = redact_by_level(run.text, redaction_level)

        output_pptx = BytesIO()
        presentation.save(output_pptx)
        output_pptx.seek(0)

        await record_history(user_id, filename, "PPTX Redaction", redaction_level, "SUCCESS")

        return StreamingResponse(
            output_pptx,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f"attachment; filename=redacted_{filename}"}
        )
    except Exception as e:
        await record_history(user_id, filename, "PPTX Redaction", redaction_level, "FAILED", str(e))
        raise HTTPException(status_code=500, detail=f"Error processing PPTX file: {str(e)}")

async def handle_csv(file_data: bytes, filename: str, redaction_level: int, user_id: Optional[int] = None) -> StreamingResponse:
    try:
        try:
            data = pd.read_csv(BytesIO(file_data), encoding="utf-8")
        except UnicodeDecodeError:
            raise HTTPException(status_code=400, detail="Error reading CSV file. Ensure UTF-8 encoding.")

        def redact_cell(cell):
            if isinstance(cell, str):
                return redact_by_level(cell, redaction_level)
            return cell

        redacted_data = data.applymap(redact_cell) if hasattr(data, 'applymap') else data.map(redact_cell)

        output_csv = BytesIO()
        redacted_data.to_csv(output_csv, index=False, encoding="utf-8")
        output_csv.seek(0)

        await record_history(user_id, filename, "CSV Redaction", redaction_level, "SUCCESS")

        return StreamingResponse(
            output_csv,
            media_type="text/csv",
            headers={"Content-Disposition": f"attachment; filename=redacted_{filename}"}
        )
    except HTTPException as http_err:
        await record_history(user_id, filename, "CSV Redaction", redaction_level, "FAILED", str(http_err.detail))
        raise http_err
    except Exception as e:
        await record_history(user_id, filename, "CSV Redaction", redaction_level, "FAILED", str(e))
        raise HTTPException(status_code=500, detail=f"Unexpected error processing CSV file: {str(e)}")

async def handle_xlsx(file_data: bytes, filename: str, redaction_level: int, user_id: Optional[int] = None) -> StreamingResponse:
    try:
        try:
            data_dict = pd.read_excel(BytesIO(file_data), sheet_name=None)
        except Exception as read_err:
            raise HTTPException(status_code=400, detail=f"Error reading Excel file: {str(read_err)}")

        def redact_cell(cell):
            if isinstance(cell, str):
                return redact_by_level(cell, redaction_level)
            return cell

        output_xlsx = BytesIO()
        with pd.ExcelWriter(output_xlsx, engine='openpyxl') as writer:
            for sheet_name, df in data_dict.items():
                redacted_df = df.applymap(redact_cell) if hasattr(df, 'applymap') else df.map(redact_cell)
                redacted_df.to_excel(writer, sheet_name=sheet_name, index=False)

        output_xlsx.seek(0)
        await record_history(user_id, filename, "XLSX Redaction", redaction_level, "SUCCESS")

        return StreamingResponse(
            output_xlsx,
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            headers={"Content-Disposition": f"attachment; filename=redacted_{filename}"}
        )
    except HTTPException as http_err:
        await record_history(user_id, filename, "XLSX Redaction", redaction_level, "FAILED", str(http_err.detail))
        raise http_err
    except Exception as e:
        await record_history(user_id, filename, "XLSX Redaction", redaction_level, "FAILED", str(e))
        raise HTTPException(status_code=500, detail=f"Unexpected error processing Excel file: {str(e)}")

async def handle_log(file_data: bytes, filename: str, redaction_level: int, user_id: Optional[int] = None) -> StreamingResponse:
    try:
        text = file_data.decode("utf-8", errors="replace")
        text = redact_by_level(text, redaction_level)
        output_file = BytesIO(text.encode("utf-8"))

        await record_history(user_id, filename, "LOG Redaction", redaction_level, "SUCCESS")

        return StreamingResponse(
            output_file,
            media_type="text/plain",
            headers={"Content-Disposition": f"attachment; filename=redacted_{filename}"}
        )
    except Exception as e:
        await record_history(user_id, filename, "LOG Redaction", redaction_level, "FAILED", str(e))
        raise HTTPException(status_code=500, detail=f"Error processing LOG file: {str(e)}")

async def handle_image(file_data: bytes, filename: str, redaction_level: int, user_id: Optional[int] = None) -> StreamingResponse:
    try:
        image = Image.open(BytesIO(file_data))
        if image.mode != 'RGB':
            image = image.convert('RGB')
        
        img_format = "JPEG" if filename.lower().endswith(('.jpg', '.jpeg')) else "PNG"
        
        if redaction_level == 0:
            output_image = BytesIO()
            image.save(output_image, format=img_format)
            output_image.seek(0)
            return StreamingResponse(
                output_image,
                media_type="image/jpeg" if img_format == "JPEG" else "image/png",
                headers={"Content-Disposition": f"attachment; filename=redacted_{filename}"}
            )

        # Automatically de-identify human faces and biometrics based on redaction level
        image = redact_faces_in_image(image, method="blur", redaction_level=redaction_level)

        draw = ImageDraw.Draw(image)
        tsv_data = pytesseract.image_to_data(image, output_type=pytesseract.Output.DATAFRAME)
        
        valid_rows = tsv_data[tsv_data['text'].notnull() & (tsv_data['text'].str.strip() != "")]
        full_image_text = " ".join(valid_rows['text'].astype(str).tolist())
        redacted_image_text = redact_by_level(full_image_text, redaction_level) if full_image_text else ""

        for _, row in valid_rows.iterrows():
            text = str(row["text"]).strip()
            if text:
                redacted_word = redact_by_level(text, redaction_level)
                if redacted_word != text or (len(text) > 2 and text not in redacted_image_text):
                    left, top, width, height = int(row['left']), int(row['top']), int(row['width']), int(row['height'])
                    draw.rectangle([left, top, left + width, top + height], fill="black")

        output_image = BytesIO()
        image.save(output_image, format=img_format)
        output_image.seek(0)

        await record_history(user_id, filename, "Image Redaction", redaction_level, "SUCCESS")

        return StreamingResponse(
            output_image,
            media_type="image/jpeg" if img_format == "JPEG" else "image/png",
            headers={"Content-Disposition": f"attachment; filename=redacted_{filename}"}
        )
    except Exception as e:
        await record_history(user_id, filename, "Image Redaction", redaction_level, "FAILED", str(e))
        raise HTTPException(status_code=500, detail=f"Error processing image: {str(e)}")
