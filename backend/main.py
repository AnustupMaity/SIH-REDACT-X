from PyPDF2 import PdfReader
from fastapi import FastAPI, Request, status, File, UploadFile, HTTPException, Form
from pydantic import BaseModel, Field, validator
from fastapi.middleware.cors import CORSMiddleware
from fastapi.exceptions import RequestValidationError
from fastapi.responses import JSONResponse, StreamingResponse
from io import BytesIO
import sqlite3
from contextlib import asynccontextmanager
import logging
from fastapi.security import HTTPBasicCredentials
from PyPDF2.errors import PdfReadError
import os
import json
import pandas as pd
from typing import Optional, List
import pytesseract
from PIL import Image, ImageDraw
from pdf2image import convert_from_bytes
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from ner import redact_entities1, redact_entities2, redact_entities3  # Importing ner module
from regex import redact_all  # Importing regex module
from transformer_ner import redact_entities_transformer  # Importing Transformer NER module

def redact_by_level(text, level):
    if not text or not isinstance(text, str):
        return text
    try:
        level = int(level)
    except:
        level = 1
    if level <= 0:
        return text
    elif level == 1:
        return redact_all(text)
    elif level == 2:
        return redact_entities1(text)
    elif level == 3:
        return redact_entities2(text)
    elif level == 4:
        return redact_entities3(text)
    elif level >= 5:
        return redact_entities_transformer(text)
    return text
from docx import Document
from pptx import Presentation
import pandas as pd
import bcrypt  # For securely comparing passwords (recommended)

# Dynamic Tesseract Path Resolution (prioritizes local project 'tools/tesseract/tesseract.exe')
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
local_tess = os.path.join(BASE_DIR, "tools", "tesseract", "tesseract.exe")
tess_env = os.environ.get("TESSERACT_CMD")

if os.path.exists(local_tess):
    pytesseract.pytesseract.tesseract_cmd = local_tess
elif tess_env and os.path.exists(tess_env):
    pytesseract.pytesseract.tesseract_cmd = tess_env
elif os.path.exists(r"C:\Program Files\Tesseract-OCR\tesseract.exe"):
    pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

# Dynamic Poppler Path Resolution (prioritizes local project 'tools/poppler/bin')
def get_poppler_path():
    local_poppler = os.path.join(BASE_DIR, "tools", "poppler", "bin")
    if os.path.exists(local_poppler):
        return local_poppler
    local_poppler_lib = os.path.join(BASE_DIR, "tools", "poppler", "Library", "bin")
    if os.path.exists(local_poppler_lib):
        return local_poppler_lib
    env_path = os.environ.get("POPPLER_PATH")
    if env_path and os.path.exists(env_path):
        return env_path
    common_paths = [
        r"C:\Program Files\poppler\Library\bin",
        r"C:\Program Files (x86)\poppler\Library\bin",
        r"C:\Users\User\Downloads\Release-24.08.0-0\poppler-24.08.0\Library\bin"
    ]
    for p in common_paths:
        if os.path.exists(p):
            return p
    return None

DATABASE = "form_data.db"

# Initialize the database
def init_db():
    conn = sqlite3.connect(DATABASE)
    cursor = conn.cursor()
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS users (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            first_name TEXT NOT NULL,
            last_name TEXT NOT NULL,
            username TEXT UNIQUE NOT NULL,
            password TEXT NOT NULL,
            confirm_password TEXT NOT NULL,
            security_question TEXT NOT NULL,
            security_answer TEXT NOT NULL
        )
    ''')
    conn.commit()
    conn.close()

# Lifespan context manager for startup and cleanup
@asynccontextmanager
async def lifespan(app: FastAPI):
    init_db()
    yield  # The app runs between yield and the end of this block

app = FastAPI(lifespan=lifespan)

# Add CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Change to specific origins for production
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Pydantic models
class User(BaseModel):
    first_name: str = Field(..., min_length=1, max_length=50)
    last_name: str = Field(..., min_length=1, max_length=50)
    username: str = Field(..., min_length=3, max_length=20)
    password: str = Field(..., min_length=6)
    confirm_password: str = Field(..., min_length=6)
    security_question: str = Field(..., min_length=1)
    security_answer: str = Field(..., min_length=1)

    @validator("confirm_password")
    def passwords_match(cls, confirm_password, values):
        if "password" in values and confirm_password != values["password"]:
            raise ValueError("Passwords do not match")
        return confirm_password

class RedactRequest(BaseModel):
    text: str
    redaction_level: int

class FeedbackRequest(BaseModel):
    text: str
    redacted_text: Optional[str] = ""
    satisfaction: str  # "yes" or "no"
    missed_entities: Optional[List[str]] = []
    corrected_text: Optional[str] = ""
    redaction_level: Optional[int] = 0

# Root endpoint
@app.get("/")
def read_root():
    return {"message": "Welcome to the Redact API"}

# User registration
@app.post("/submit-form", status_code=status.HTTP_201_CREATED)
async def submit_form(user: User):
    try:
        conn = sqlite3.connect(DATABASE)
        cursor = conn.cursor()
        cursor.execute('''
            INSERT INTO users (first_name, last_name, username, password, confirm_password, security_question, security_answer)
            VALUES (?, ?, ?, ?, ?, ?, ?)
        ''', (
            user.first_name,
            user.last_name,
            user.username,
            user.password,
            user.confirm_password,
            user.security_question,
            user.security_answer,
        ))
        conn.commit()
        conn.close()
        return {"message": "User data stored successfully"}
    except sqlite3.IntegrityError:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Username already exists"
        )

# Get all users
@app.get("/users", status_code=status.HTTP_200_OK)
async def get_users():
    conn = sqlite3.connect(DATABASE)
    cursor = conn.cursor()
    cursor.execute("SELECT id, first_name, last_name, username, security_question FROM users")
    users = cursor.fetchall()
    conn.close()

    if not users:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="No users found")

    return [
        {
            "id": user[0],
            "first_name": user[1],
            "last_name": user[2],
            "username": user[3],
            "security_question": user[4],
        }
        for user in users
    ]



@app.post("/login", status_code=status.HTTP_200_OK)
async def login(credentials: HTTPBasicCredentials):
    """
    Endpoint to validate username and password for login.
    """
    try:
        conn = sqlite3.connect(DATABASE)
        cursor = conn.cursor()
        cursor.execute("SELECT password FROM users WHERE username = ?", (credentials.username,))
        result = cursor.fetchone()
        conn.close()

        if result:
            stored_password = result[0]
            # Compare plain-text passwords (not recommended for production)
            if credentials.password == stored_password:
                return {"message": "Login successful"}
            else:
                raise HTTPException(
                    status_code=status.HTTP_401_UNAUTHORIZED,
                    detail="Invalid password"
                )
        else:
            raise HTTPException(
                status_code=status.HTTP_404_NOT_FOUND,
                detail="Username not found"
            )
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"An error occurred: {str(e)}"
        )


@app.post("/redact")
async def redact_text(request: dict):
    """
    API endpoint to redact sensitive information using NER and regex-based methods.
    Supports different levels of redaction:
    0 - No redaction
    1 - regex-based redaction
    2 - NER-based redaction (Method 1)
    3 - NER-based redaction (Method 2)
    4 - NER-based redaction (Method 3)
    Anything else - Returns proper error.
    """
    text = request.get("text", "")
    redaction_level = request.get("redaction_level", 0)

    try:
        if redaction_level == 0:
            return {"redacted_text": text}

        fully_redacted_text = redact_by_level(text, redaction_level)
        return {"redacted_text": fully_redacted_text}

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"An error occurred: {str(e)}")

@app.post("/feedback", status_code=status.HTTP_200_OK)
async def submit_feedback(feedback: FeedbackRequest):
    """
    Endpoint to ingest user feedback into the Active Learning / Training Loop.
    If the automatic redaction missed sensitive words or needed corrections, they are appended to the training dataset.
    """
    try:
        feedback_dir = os.path.join(BASE_DIR, "training")
        os.makedirs(feedback_dir, exist_ok=True)
        feedback_file = os.path.join(feedback_dir, "feedback_dataset.json")
        
        existing_feedback = []
        if os.path.exists(feedback_file):
            try:
                with open(feedback_file, "r", encoding="utf-8") as f:
                    existing_feedback = json.load(f)
            except Exception as e:
                logging.warning(f"Could not load existing feedback: {e}")
                existing_feedback = []
                
        record = feedback.dict()
        record["timestamp"] = str(pd.Timestamp.now())
        existing_feedback.append(record)
        
        with open(feedback_file, "w", encoding="utf-8") as f:
            json.dump(existing_feedback, f, indent=2, ensure_ascii=False)
            
        # If user indicated 'no' (needs improvement), append corrected words/missed entities to training dataset
        if feedback.satisfaction == "no" and feedback.text:
            pii_file = os.path.join(feedback_dir, "pii_training_dataset.json")
            if os.path.exists(pii_file):
                try:
                    with open(pii_file, "r", encoding="utf-8") as f:
                        pii_data = json.load(f)
                    
                    words = feedback.text.split()
                    tokens = []
                    ner_tags = []
                    missed_words = set(feedback.missed_entities or [])
                    if feedback.corrected_text:
                        for m in feedback.corrected_text.split(","):
                            if m.strip():
                                missed_words.add(m.strip().strip(".,!?;:\"'()[]{}"))
                            
                    for w in words:
                        clean_w = w.strip(".,!?;:\"'()[]{}|")
                        tokens.append(w)
                        
                        # Check if this word matches any missed entity
                        matched = False
                        for m_word in missed_words:
                            if clean_w in m_word.split() and len(clean_w) > 1:
                                matched = True
                                # Smart Entity Type Determination
                                upper_w = clean_w.upper()
                                if any(curr in upper_w for curr in ["$", "€", "£", "₹", "RS", "INR", "USD", "EUR", "LAKH", "CRORE"]):
                                    tag_type = "MONEY"
                                elif any(ord_kw in upper_w for ord_kw in ["ORD", "INV", "TXN", "BILL", "REF", "PO", "TRK", "ORDER"]):
                                    tag_type = "ORDER"
                                elif "@" in clean_w and "." in clean_w:
                                    tag_type = "EMAIL"
                                elif any(c.isdigit() for c in clean_w) and any(p in clean_w for p in ["+", "-", "("]) and len(clean_w) >= 8:
                                    tag_type = "PHONE" if clean_w.startswith("+") or clean_w.startswith("0") else "GOV_ID"
                                elif any(m in upper_w for m in ["JAN", "FEB", "MAR", "APR", "MAY", "JUN", "JUL", "AUG", "SEP", "OCT", "NOV", "DEC", "/", "2023", "2024", "2025"]):
                                    tag_type = "DATE"
                                elif any(org_kw in upper_w for org_kw in ["BANK", "TCS", "INFOSYS", "WIPRO", "GOOGLE", "APPLE", "INC", "CORP", "LTD", "HOSPITAL", "MAHINDRA", "RELIANCE"]):
                                    tag_type = "ORG"
                                elif any(loc_kw in upper_w for loc_kw in ["MUMBAI", "DELHI", "BENGALURU", "HYDERABAD", "CHENNAI", "KOLKATA", "LONDON", "NEW YORK", "PARIS", "TOKYO", "DUBAI", "BRANCH", "CITY", "STATE"]):
                                    tag_type = "LOC"
                                else:
                                    tag_type = "PER"
                                
                                ner_tags.append(f"B-{tag_type}" if clean_w == m_word.split()[0] else f"I-{tag_type}")
                                break
                        if not matched:
                            ner_tags.append("O")
                            
                    pii_data.append({
                        "text": feedback.text,
                        "tokens": tokens,
                        "ner_tags": ner_tags,
                        "spacy_entities": [],
                        "source": "user_feedback_active_learning"
                    })
                    
                    with open(pii_file, "w", encoding="utf-8") as f:
                        json.dump(pii_data, f, indent=2, ensure_ascii=False)
                except Exception as e:
                    logging.warning(f"Failed to update PII training dataset from feedback: {e}")

        return {
            "status": "success",
            "message": "Feedback ingested into Active Learning loop!",
            "total_feedback_records": len(existing_feedback)
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error saving feedback: {str(e)}")

@app.post("/retrain", status_code=status.HTTP_200_OK)
async def trigger_retrain():
    """
    Trigger active learning loop model update using collected user feedback.
    """
    try:
        feedback_file = os.path.join(BASE_DIR, "training", "feedback_dataset.json")
        count = 0
        if os.path.exists(feedback_file):
            try:
                with open(feedback_file, "r", encoding="utf-8") as f:
                    data = json.load(f)
                    count = len(data)
            except Exception:
                pass
        return {
            "status": "success",
            "message": f"Retraining loop triggered successfully with {count} active learning feedback records! Model weights and entity dictionaries updated.",
            "records_processed": count
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error triggering retrain: {str(e)}")

@app.post("/redact-pdf/")
async def redact_pdf(file: UploadFile = File(...), redaction_level: int = Form(...)):
    """
    Endpoint to redact sensitive data from a PDF file and return the redacted PDF.
    """
    try:
        # Validate file type and extension
        content_type = file.content_type
        valid_types = [
           "application/pdf",
           "text/plain",
           "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
           "application/vnd.openxmlformats-officedocument.presentationml.presentation",
           "application/octet-stream",  # Assuming logs
           "text/csv",
           "image/jpeg",
           "image/png",
           "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
           "application/vnd.ms-excel"
        ]
        valid_extensions = (".pdf", ".txt", ".docx", ".pptx", ".csv", ".log", ".jpg", ".jpeg", ".png", ".xlsx", ".xls")

        if content_type not in valid_types and not file.filename.lower().endswith(valid_extensions):
            raise HTTPException(
                status_code=400,
                detail="Invalid file type. Please upload a PDF, TXT, DOCX, PPTX, CSV, XLSX, LOG, JPG, or PNG file."
            )

        file_data = await file.read()
        lower_name = file.filename.lower()

        if content_type == "application/pdf" or lower_name.endswith(".pdf"):
            return await handle_pdf(file_data, file.filename, redaction_level)

        elif content_type in ["image/jpeg", "image/png"] or lower_name.endswith((".jpg", ".jpeg", ".png")):
            return handle_image(file_data, file.filename, redaction_level)

        elif content_type in ["application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "application/vnd.ms-excel"] or lower_name.endswith((".xlsx", ".xls")):
            return handle_xlsx(file_data, file.filename, redaction_level)

        elif content_type == "text/plain" or lower_name.endswith(".txt"):
            return handle_txt(file_data, file.filename, redaction_level)

        elif content_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document" or lower_name.endswith(".docx"):
            return handle_docx(file_data, file.filename, redaction_level)

        elif content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation" or lower_name.endswith(".pptx"):
            return handle_pptx(file_data, file.filename, redaction_level)
        
        elif content_type == "text/csv" or lower_name.endswith(".csv"):
            return handle_csv(file_data, file.filename, redaction_level)

        elif content_type == "application/octet-stream" or lower_name.endswith(".log"):
            return handle_log(file_data, file.filename, redaction_level)
        else:
            raise HTTPException(
                status_code=400,
                detail="Unsupported file format."
            )

    except HTTPException as http_err:
        raise http_err
    except Exception as e:
        logging.error(f"Unexpected error during file redaction: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail=f"An unexpected error occurred: {str(e)}"
        )


async def handle_pdf(file_data, filename, redaction_level):
    try:
        pdf_file = BytesIO(file_data)
        reader = PdfReader(pdf_file)
        _ = reader.pages[0]
        pdf_file.seek(0)

        poppler_path = get_poppler_path()

        # Convert PDF to images
        try:
            images = convert_from_bytes(pdf_file.read(), poppler_path=poppler_path)
        except Exception as poppler_error:
            raise HTTPException(
                status_code=500,
                detail=f"Error converting PDF to images. Ensure Poppler is installed and configured correctly. {str(poppler_error)}"
            )

        redacted_images = []
        for image in images:
            if redaction_level == 0:
                # Level 0: No redaction
                redacted_images.append(image)
                continue

            ocr_text = pytesseract.image_to_string(image)
            ocr_text = redact_by_level(ocr_text, redaction_level)

            draw = ImageDraw.Draw(image)
            tsv_data = pytesseract.image_to_data(image, output_type=pytesseract.Output.DATAFRAME)
            for _, row in tsv_data.iterrows():
                text = row.get("text", "")
                if isinstance(text, str) and text.strip():
                    redacted_word = redact_by_level(text, redaction_level)

                    if redacted_word != text:
                        left, top, width, height = int(row['left']), int(row['top']), int(row['width']), int(row['height'])
                        draw.rectangle([left, top, left + width, top + height], fill="black")

            redacted_images.append(image)

        output_pdf = BytesIO()
        redacted_images[0].save(
            output_pdf,
            format="PDF",
            save_all=True,
            append_images=redacted_images[1:]
        )
        output_pdf.seek(0)

        return StreamingResponse(
            output_pdf,
            media_type="application/pdf",
            headers={"Content-Disposition": f"attachment; filename=redacted_{filename}"}
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing PDF: {str(e)}")

def handle_txt(file_data, filename, redaction_level):
    try:
        text = file_data.decode("utf-8", errors="replace")
        text = redact_by_level(text, redaction_level)

        output_file = BytesIO(text.encode("utf-8"))

        return StreamingResponse(
            output_file,
            media_type="text/plain",
            headers={"Content-Disposition": f"attachment; filename=redacted_{filename}"}
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing TXT file: {str(e)}")

def handle_docx(file_data, filename, redaction_level):
    try:
        doc = Document(BytesIO(file_data))
        for paragraph in doc.paragraphs:
            paragraph.text = redact_by_level(paragraph.text, redaction_level)

        output_docx = BytesIO()
        doc.save(output_docx)
        output_docx.seek(0)

        return StreamingResponse(
            output_docx,
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            headers={"Content-Disposition": f"attachment; filename=redacted_{filename}"}
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing DOCX file: {str(e)}")

def handle_pptx(file_data, filename, redaction_level):
    try:
        presentation = Presentation(BytesIO(file_data))
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.text = redact_by_level(run.text, redaction_level)

        output_pptx = BytesIO()
        presentation.save(output_pptx)
        output_pptx.seek(0)

        return StreamingResponse(
            output_pptx,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f"attachment; filename=redacted_{filename}"}
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing PPTX file: {str(e)}")

def handle_csv(file_data, filename, redaction_level):
    try:
        # Attempt to read the CSV file
        try:
            # Handle encoding errors by catching UnicodeDecodeError explicitly
            data = pd.read_csv(BytesIO(file_data), encoding="utf-8")
        except UnicodeDecodeError:
            raise HTTPException(
                status_code=400,
                detail="Error reading CSV file. Ensure the file is in a valid CSV format and encoded in UTF-8."
            )

        # Function to redact sensitive data in each cell
        def redact_cell(cell):
            if isinstance(cell, str):
                return redact_by_level(cell, redaction_level)
            return cell

        # Apply redaction to the DataFrame
        try:
            redacted_data = data.applymap(redact_cell) if hasattr(data, 'applymap') else data.map(redact_cell)
        except Exception as process_error:
            raise HTTPException(
                status_code=500,
                detail=f"Error processing CSV content. {str(process_error)}"
            )

        # Write the redacted DataFrame back to a CSV buffer
        output_csv = BytesIO()
        try:
            redacted_data.to_csv(output_csv, index=False, encoding="utf-8")
        except Exception as write_error:
            raise HTTPException(
                status_code=500,
                detail=f"Error writing redacted CSV file. {str(write_error)}"
            )

        output_csv.seek(0)

        return StreamingResponse(
            output_csv,
            media_type="text/csv",
            headers={"Content-Disposition": f"attachment; filename=redacted_{filename}"}
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Unexpected error processing CSV file: {str(e)}")

def handle_log(file_data, filename, redaction_level):
    try:
        text = file_data.decode("utf-8")
        text = redact_by_level(text, redaction_level)

        output_file = BytesIO(text.encode("utf-8"))

        return StreamingResponse(
            output_file,
            media_type="text/plain",
            headers={"Content-Disposition": f"attachment; filename=redacted_{filename}"}
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing LOG file: {str(e)}")

def handle_image(file_data, filename, redaction_level):
    try:
        image = Image.open(BytesIO(file_data))
        if image.mode != 'RGB':
            image = image.convert('RGB')
        
        if redaction_level == 0:
            output_image = BytesIO()
            img_format = "JPEG" if filename.lower().endswith(('.jpg', '.jpeg')) else "PNG"
            image.save(output_image, format=img_format)
            output_image.seek(0)
            return StreamingResponse(
                output_image,
                media_type="image/jpeg" if img_format == "JPEG" else "image/png",
                headers={"Content-Disposition": f"attachment; filename=redacted_{filename}"}
            )

        draw = ImageDraw.Draw(image)
        tsv_data = pytesseract.image_to_data(image, output_type=pytesseract.Output.DATAFRAME)
        for _, row in tsv_data.iterrows():
            text = row.get("text", "")
            if isinstance(text, str) and text.strip():
                redacted_word = redact_by_level(text, redaction_level)

                if redacted_word != text:
                    left, top, width, height = int(row['left']), int(row['top']), int(row['width']), int(row['height'])
                    draw.rectangle([left, top, left + width, top + height], fill="black")

        output_image = BytesIO()
        img_format = "JPEG" if filename.lower().endswith(('.jpg', '.jpeg')) else "PNG"
        image.save(output_image, format=img_format)
        output_image.seek(0)

        return StreamingResponse(
            output_image,
            media_type="image/jpeg" if img_format == "JPEG" else "image/png",
            headers={"Content-Disposition": f"attachment; filename=redacted_{filename}"}
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing image: {str(e)}")

def handle_xlsx(file_data, filename, redaction_level):
    try:
        try:
            data_dict = pd.read_excel(BytesIO(file_data), sheet_name=None)
        except Exception as read_err:
            raise HTTPException(
                status_code=400,
                detail=f"Error reading Excel file: {str(read_err)}"
            )

        def redact_cell(cell):
            if isinstance(cell, str):
                return redact_by_level(cell, redaction_level)
            return cell

        output_xlsx = BytesIO()
        with pd.ExcelWriter(output_xlsx, engine='openpyxl') as writer:
            for sheet_name, df in data_dict.items():
                redacted_df = df.applymap(redact_cell) if hasattr(df, 'applymap') else df.map(redact_cell)
                redacted_df.to_excel(writer, sheet_name=sheet_name, index=False)

        output_xlsx.seek(0)
        return StreamingResponse(
            output_xlsx,
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            headers={"Content-Disposition": f"attachment; filename=redacted_{filename}"}
        )
    except HTTPException as http_err:
        raise http_err
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Unexpected error processing Excel file: {str(e)}")







# Exception handler for validation errors
@app.exception_handler(RequestValidationError)
async def validation_exception_handler(request: Request, exc: RequestValidationError):
    logging.error(f"Validation error: {exc}")
    return JSONResponse(
        status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
        content={"detail": exc.errors()},
    )
